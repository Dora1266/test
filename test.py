import os
import io
import re
import fitz
import docx
import csv
import json
import uuid
import time
import logging
import tempfile
import argparse
import requests
import concurrent.futures
import nltk
from nltk.tokenize import sent_tokenize
from typing import List, Dict, Any, Tuple
from flask import Flask, request, jsonify, render_template, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger('document_dialogue_app')

# Try to import optional dependencies
try:
    import spacy
    import easyocr
    import numpy as np
    from PIL import Image
    from pptx import Presentation
    from openpyxl import load_workbook

    # Initialize spaCy
    try:
        nlp = spacy.load('zh_core_web_sm')
    except OSError:
        try:
            nlp = spacy.load('en_core_web_sm')
        except OSError:
            nlp = None

    # Initialize NLTK
    try:
        nltk.data.find('tokenizers/punkt')
    except LookupError:
        nltk.download('punkt', quiet=True)

    # Flag to indicate optional dependencies are available
    OPTIONAL_DEPS_AVAILABLE = True
except ImportError:
    OPTIONAL_DEPS_AVAILABLE = False
    nlp = None
    logger.warning("Some optional dependencies are missing. Limited functionality available.")


#################################################
# Document Parsing Component
#################################################

class SemanticDocumentParser:
    def __init__(self, language_list=None, segment_size=1000, overlap=True, overlap_limit=200,
                 clean_for_ai=True, segmentation_type="semantic", replace_whitespace=False,
                 remove_urls_emails=False):
        if language_list is None:
            language_list = ['ch_sim', 'en']

        self.languages = language_list
        self.segment_size = segment_size
        self.overlap = overlap
        self.overlap_limit = overlap_limit
        self.nlp = nlp
        self.clean_for_ai = clean_for_ai
        self.segmentation_type = segmentation_type
        self.replace_whitespace = replace_whitespace
        self.remove_urls_emails = remove_urls_emails

        # Language detection
        self.primary_language = 'en'
        if 'ch_sim' in language_list or 'ch_tra' in language_list:
            self.primary_language = 'zh'

        # Initialize OCR if available
        if OPTIONAL_DEPS_AVAILABLE:
            try:
                self.reader = easyocr.Reader(language_list)
            except Exception as e:
                logger.error(f"Failed to initialize EasyOCR: {e}")
                self.reader = None
        else:
            self.reader = None

    def parse(self, file_path, segment_size=None):
        if segment_size is not None:
            self.segment_size = segment_size

        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        if ext == '.pdf':
            return self.parse_pdf(file_path)
        elif ext == '.docx':
            return self.parse_docx(file_path)
        elif ext == '.txt':
            return self.parse_txt(file_path)
        elif ext == '.pptx':
            return self.parse_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self.parse_excel(file_path)
        elif ext in ['.csv']:
            return self.parse_csv(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def parse_pdf(self, file_path):
        if not OPTIONAL_DEPS_AVAILABLE:
            raise ValueError("PDF parsing requires additional dependencies")

        doc = fitz.open(file_path)
        full_text = []
        page_texts = []  # For structural segmentation

        # First, extract all text content from PDF
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text()

            # If page has very little text, try OCR
            if len(page_text.strip()) < 50 and self.reader:
                pix = page.get_pixmap()
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                img_np = np.array(img)
                ocr_results = self.reader.readtext(img_np)
                page_text_from_ocr = ' '.join([text for _, text, _ in ocr_results])

                if page_text_from_ocr.strip():
                    page_text = page_text_from_ocr

            if self.segmentation_type == "structural":
                # Store each page as a separate segment
                if page_text.strip():
                    page_texts.append({
                        'content': self.clean_text(page_text),
                        'type': 'page',
                        'page_num': page_num + 1
                    })
            else:
                full_text.append(page_text)

            # Process images in the page
            if self.reader:
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]

                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        img_np = np.array(img)
                        ocr_results = self.reader.readtext(img_np)
                        image_text = ' '.join([text for _, text, _ in ocr_results])

                        if image_text.strip():
                            if self.segmentation_type == "structural":
                                page_texts.append({
                                    'content': self.clean_text(image_text),
                                    'type': 'image',
                                    'page_num': page_num + 1,
                                    'image_num': img_index + 1
                                })
                            else:
                                full_text.append(image_text)
                    except Exception as e:
                        logger.error(f"Error processing image: {e}")

        if self.segmentation_type == "structural":
            return self.apply_segment_size_limit(page_texts)
        else:
            # Join all text and process by semantic units
            combined_text = '\n'.join(full_text)
            return self.segment_by_semantics(combined_text)

    def parse_docx(self, file_path):
        doc = docx.Document(file_path)
        paragraphs_text = []

        # Extract paragraphs, which naturally tend to be semantic units
        for para_index, para in enumerate(doc.paragraphs):
            if para.text.strip():
                cleaned_text = self.clean_text(para.text)
                if self.segmentation_type == "structural":
                    paragraphs_text.append({
                        'content': cleaned_text,
                        'type': 'paragraph',
                        'paragraph_num': para_index + 1
                    })
                else:
                    paragraphs_text.append(cleaned_text)

        # Process tables
        table_texts = []
        for table_index, table in enumerate(doc.tables):
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(' | '.join(row_text))

            if table_text:
                table_content = '\n'.join(table_text)
                cleaned_table = self.clean_text(table_content)
                if self.segmentation_type == "structural":
                    table_texts.append({
                        'content': cleaned_table,
                        'type': 'table',
                        'table_num': table_index + 1
                    })
                else:
                    table_texts.append(cleaned_table)

        if self.segmentation_type == "structural":
            # Return structural segments
            all_segments = paragraphs_text + table_texts
            return self.apply_segment_size_limit(all_segments)
        else:
            # Combine all text and process semantically
            combined_text = '\n\n'.join(paragraphs_text + table_texts)
            return self.segment_by_semantics(combined_text)

    def parse_txt(self, file_path):
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read()

        if self.segmentation_type == "structural":
            # Split by paragraphs (double newlines)
            paragraphs = re.split(r'\n\s*\n', text)
            segments = []
            for i, para in enumerate(paragraphs):
                if para.strip():
                    segments.append({
                        'content': self.clean_text(para),
                        'type': 'paragraph',
                        'paragraph_num': i + 1
                    })
            return self.apply_segment_size_limit(segments)
        else:
            return self.segment_by_semantics(text)

    def parse_pptx(self, file_path):
        if not OPTIONAL_DEPS_AVAILABLE:
            raise ValueError("PowerPoint parsing requires additional dependencies")

        prs = Presentation(file_path)
        slide_texts = []

        for slide_index, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)

            if texts:
                slide_content = '\n'.join(texts)
                cleaned_slide = self.clean_text(slide_content)
                if self.segmentation_type == "structural":
                    slide_texts.append({
                        'content': cleaned_slide,
                        'type': 'slide',
                        'slide_num': slide_index + 1
                    })
                else:
                    slide_texts.append(cleaned_slide)

        if self.segmentation_type == "structural":
            return self.apply_segment_size_limit(slide_texts)
        else:
            # Join with double newlines to indicate strong semantic breaks between slides
            combined_text = '\n\n'.join(slide_texts)
            return self.segment_by_semantics(combined_text)

    def parse_excel(self, file_path):
        if not OPTIONAL_DEPS_AVAILABLE:
            raise ValueError("Excel parsing requires additional dependencies")

        wb = load_workbook(file_path, data_only=True)
        sheet_texts = []

        for sheet_index, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            rows = []

            for row in sheet.iter_rows(values_only=True):
                if any(cell for cell in row):
                    row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                    rows.append(row_text)

            if rows:
                sheet_content = '\n'.join(rows)
                cleaned_sheet = self.clean_text(sheet_content)
                if self.segmentation_type == "structural":
                    sheet_texts.append({
                        'content': cleaned_sheet,
                        'type': 'sheet',
                        'sheet_name': sheet_name,
                        'sheet_index': sheet_index + 1
                    })
                else:
                    sheet_texts.append(cleaned_sheet)

        if self.segmentation_type == "structural":
            return self.apply_segment_size_limit(sheet_texts)
        else:
            # Each sheet is treated as a potential semantic unit
            combined_text = '\n\n'.join(sheet_texts)
            return self.segment_by_semantics(combined_text)

    def parse_csv(self, file_path):
        rows = []

        with open(file_path, 'r', newline='', encoding='utf-8', errors='ignore') as csvfile:
            csv_reader = csv.reader(csvfile)
            for row in csv_reader:
                if any(cell.strip() for cell in row):
                    rows.append(' | '.join(row))

        if rows:
            combined_text = '\n'.join(rows)
            if self.segmentation_type == "structural":
                # Treat the entire CSV as one segment or split by chunks
                chunks = []
                chunk_size = 50  # Number of rows per chunk
                for i in range(0, len(rows), chunk_size):
                    chunk = rows[i:i + chunk_size]
                    chunk_text = '\n'.join(chunk)
                    chunks.append({
                        'content': self.clean_text(chunk_text),
                        'type': 'csv_chunk',
                        'row_start': i + 1,
                        'row_end': min(i + chunk_size, len(rows))
                    })
                return self.apply_segment_size_limit(chunks)
            else:
                return self.segment_by_semantics(combined_text)
        return []

    def clean_text(self, text):
        """
        Apply basic text cleaning operations based on user preferences.
        """
        if not text or not text.strip():
            return text

        # Clean the text if needed
        if self.replace_whitespace:
            # Replace consecutive whitespace (spaces, tabs, newlines) with a single space
            text = re.sub(r'\s+', ' ', text)

        if self.remove_urls_emails:
            # Remove URLs
            text = re.sub(r'https?://\S+|www\.\S+', '', text)
            # Remove emails
            text = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '', text)

        # Apply AI cleaning if enabled
        if self.clean_for_ai:
            text = self.clean_text_for_ai_training(text)

        return text.strip()

    def identify_semantic_blocks(self, text):
        """
        Identify semantic blocks in the text using various heuristics.
        Returns a list of blocks, where each block is a list of sentences.
        """
        if not text or not text.strip():
            return []

        # Clean the text
        text = re.sub(r'[ \t]+', ' ', text).strip()
        text = re.sub(r'\n{3,}', '\n\n', text)  # Normalize multiple newlines

        # First split by obvious semantic breaks (multiple newlines)
        # These often indicate new sections or major topic shifts
        paragraph_blocks = re.split(r'\n\s*\n', text)

        semantic_blocks = []

        for block in paragraph_blocks:
            if not block.strip():
                continue

            # Further analyze the block
            if self.nlp and self.primary_language in ['en', 'zh']:
                # Use spaCy for more advanced semantic analysis if available
                doc = self.nlp(block)

                # We'll use a simple approach here, but could be extended with
                # more sophisticated analysis of entity transitions, topic modeling, etc.
                current_sentences = []
                sentence_count = 0

                for sent in doc.sents:
                    sentence_count += 1
                    current_sentences.append(sent.text)

                    # Check if we should create a new block
                    if sentence_count >= 5:  # A reasonable number for semantic cohesion
                        semantic_blocks.append(current_sentences)
                        current_sentences = []
                        sentence_count = 0

                if current_sentences:
                    semantic_blocks.append(current_sentences)
            else:
                # Fallback to simple sentence tokenization
                if self.primary_language == 'zh':
                    # For Chinese, use regex-based splitting
                    sentences = re.split(r'([。！？；])', block)
                    processed_sentences = []

                    i = 0
                    while i < len(sentences):
                        if i + 1 < len(sentences) and re.match(r'[。！？；]', sentences[i + 1]):
                            processed_sentences.append(sentences[i] + sentences[i + 1])
                            i += 2
                        else:
                            if sentences[i].strip():
                                processed_sentences.append(sentences[i])
                            i += 1
                else:
                    # For other languages, use NLTK's sentence tokenizer
                    processed_sentences = sent_tokenize(block)

                # Group sentences into semantic blocks
                current_block = []
                for sent in processed_sentences:
                    current_block.append(sent)
                    if len(current_block) >= 5:  # Same heuristic as above
                        semantic_blocks.append(current_block)
                        current_block = []

                if current_block:
                    semantic_blocks.append(current_block)

        return semantic_blocks

    def clean_text_for_ai_training(self, text):
        """
        Remove content that could negatively impact AI training, such as URLs,
        email addresses, phone numbers, and other sensitive patterns.
        """
        if not text or not text.strip():
            return text

        # Remove URLs completely
        text = re.sub(r'https?://\S+|www\.\S+', '', text)

        # Remove email addresses
        text = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '', text)

        # Remove IP addresses
        text = re.sub(r'\b(?:\d{1,3}\.){3}\d{1,3}\b', '', text)

        # Remove phone numbers - multiple formats
        text = re.sub(r'\+?\d{1,4}?[-.\s]?\(?\d{1,3}?\)?[-.\s]?\d{1,4}[-.\s]?\d{1,4}[-.\s]?\d{1,9}', '', text)

        # Remove social media handles
        text = re.sub(r'@\w+', '', text)

        # Remove credit card numbers
        text = re.sub(r'\b(?:\d{4}[-\s]?){3}\d{4}\b', '', text)

        # Remove social security numbers (US format)
        text = re.sub(r'\b\d{3}-\d{2}-\d{4}\b', '', text)

        # Remove Chinese ID numbers
        text = re.sub(r'\b\d{17}[\dXx]\b', '', text)

        # Remove API keys and tokens (common patterns)
        text = re.sub(r'\b[A-Za-z0-9_\-]{20,}\b', '', text)

        # Remove markdown URLs
        text = re.sub(r'\[.*?\]\(.*?\)', '', text)

        # Remove HTML/XML tags
        text = re.sub(r'<[^>]+>', '', text)

        # Remove excessive whitespace and normalize
        text = re.sub(r'\s+', ' ', text)

        # Fix up any double/triple/etc spaces that might have been created
        while '  ' in text:
            text = text.replace('  ', ' ')

        return text.strip()

    def segment_by_semantics(self, text):
        """
        Segment text into semantic units, then ensure each segment is within size limits.
        """
        if not text or not text.strip():
            return []

        # Apply basic cleaning according to preferences
        text = self.clean_text(text)

        semantic_blocks = self.identify_semantic_blocks(text)
        segments = []
        last_part = None

        # First pass: combine sentences within semantic blocks
        for block in semantic_blocks:
            block_text = ' '.join(block)

            # If block is too large, we need to split it further
            if len(block_text) > self.segment_size:
                # Try to split at semantically meaningful places
                sub_segments = self.split_large_block(block)

                for segment in sub_segments:
                    if last_part and self.overlap:
                        segment = self.apply_overlap(segment, last_part)

                    segments.append({
                        'content': segment,
                        'type': 'semantic_unit'
                    })
                    last_part = segment
            else:
                if last_part and self.overlap:
                    block_text = self.apply_overlap(block_text, last_part)

                segments.append({
                    'content': block_text,
                    'type': 'semantic_unit'
                })
                last_part = block_text

        # Second pass: check if any adjacent small segments can be combined
        optimized_segments = self.optimize_segment_sizes(segments)
        return optimized_segments

    def apply_segment_size_limit(self, segments):
        """
        Apply segment size limit to structural segments, splitting if necessary.
        """
        result = []
        for segment in segments:
            content = segment['content']
            if len(content) <= self.segment_size:
                result.append(segment)
            else:
                # Split the content and create multiple segments
                if self.primary_language == 'zh':
                    sentences = self.split_chinese_sentences(content)
                else:
                    sentences = sent_tokenize(content)

                # Group sentences into segments
                current_segment = ""
                segment_index = 1
                for sent in sentences:
                    if len(current_segment) + len(sent) + 1 <= self.segment_size:
                        if current_segment:
                            current_segment += " " + sent
                        else:
                            current_segment = sent
                    else:
                        # Create a new segment
                        segment_copy = segment.copy()
                        segment_copy['content'] = current_segment
                        segment_copy['part'] = segment_index
                        result.append(segment_copy)

                        current_segment = sent
                        segment_index += 1

                if current_segment:
                    segment_copy = segment.copy()
                    segment_copy['content'] = current_segment
                    segment_copy['part'] = segment_index
                    result.append(segment_copy)

        return result

    def split_chinese_sentences(self, text):
        """Split Chinese text into sentences."""
        sentences = re.split(r'([。！？；])', text)
        processed_sentences = []

        i = 0
        while i < len(sentences):
            if i + 1 < len(sentences) and re.match(r'[。！？；]', sentences[i + 1]):
                processed_sentences.append(sentences[i] + sentences[i + 1])
                i += 2
            else:
                if sentences[i].strip():
                    processed_sentences.append(sentences[i])
                i += 1

        return processed_sentences

    def split_large_block(self, sentences):
        """
        Split a large block of sentences into smaller segments, preserving semantic meaning where possible.
        """
        segments = []
        current_segment = []
        current_length = 0

        for sentence in sentences:
            # If a single sentence exceeds the limit, we'll have to split it
            if len(sentence) > self.segment_size:
                if current_segment:
                    segments.append(' '.join(current_segment))
                    current_segment = []
                    current_length = 0

                # Split long sentence by clauses or simply by size if needed
                split_parts = self.split_long_sentence(sentence)
                segments.extend(split_parts)
                continue

            # Add the sentence if it fits
            if current_length + len(sentence) + 1 <= self.segment_size:
                current_segment.append(sentence)
                current_length += len(sentence) + 1
            else:
                # This sentence would make the segment too large, so finalize current segment
                if current_segment:
                    segments.append(' '.join(current_segment))

                # Start a new segment with this sentence
                current_segment = [sentence]
                current_length = len(sentence)

        # Add any remaining content
        if current_segment:
            segments.append(' '.join(current_segment))

        return segments

    def split_long_sentence(self, sentence):
        """
        Split an exceptionally long sentence into smaller parts.
        Attempts to split at clause boundaries if possible.
        """
        # Try to split at clause boundaries first
        clause_markers = [
            ',', '，', ';', '；', ':', '：',
            'and', 'but', 'or', 'nor', 'yet', 'so',
            '和', '或者', '但是', '而且', '然后', '因此'
        ]

        parts = []
        current_part = ""
        words = re.split(r'(\s+)', sentence) if self.primary_language != 'zh' else list(sentence)

        for word in words:
            if current_part and len(current_part + word) > self.segment_size:
                parts.append(current_part.strip())
                current_part = word
            else:
                current_part += word

                # Check if we have a good breaking point
                for marker in clause_markers:
                    if word.endswith(marker) and len(current_part) > self.segment_size // 2:
                        parts.append(current_part.strip())
                        current_part = ""
                        break

        if current_part:
            parts.append(current_part.strip())

        # If we still have parts that are too long, just use brute force
        final_parts = []
        for part in parts:
            if len(part) <= self.segment_size:
                final_parts.append(part)
            else:
                # Just split by size
                for i in range(0, len(part), self.segment_size):
                    final_parts.append(part[i:i + self.segment_size])

        return final_parts

    def optimize_segment_sizes(self, segments):
        """
        Try to optimize segment sizes by combining small adjacent segments
        that together stay within the size limit.
        """
        if not segments or len(segments) < 2:
            return segments

        optimized = []
        i = 0

        while i < len(segments):
            current = segments[i]

            # Look ahead to see if we can combine with the next segment
            if (i + 1 < len(segments) and
                    len(current['content']) + len(segments[i + 1]['content']) + 1 <= self.segment_size):

                # Combine the segments
                combined_content = current['content'] + ' ' + segments[i + 1]['content']
                optimized.append({
                    'content': combined_content,
                    'type': 'semantic_unit'
                })
                i += 2
            else:
                optimized.append(current)
                i += 1

        return optimized

    def apply_overlap(self, current_text, last_part):
        """
        Apply overlap between segments to maintain context.
        Now only overlaps the last sentence from the previous segment.
        """
        if not current_text.strip() or not self.overlap or not last_part:
            return current_text

        # Get the last sentence of the previous part
        if self.nlp:
            doc = self.nlp(last_part)
            sentences = list(doc.sents)
            if sentences:
                last_sentence = sentences[-1].text
            else:
                return current_text
        else:
            # Fallback method to get last sentence
            if self.primary_language == 'zh':
                sentences = self.split_chinese_sentences(last_part)
            else:
                sentences = sent_tokenize(last_part)

            if sentences:
                last_sentence = sentences[-1]
            else:
                return current_text

        # Trim if too long
        if len(last_sentence) > self.overlap_limit:
            last_sentence = last_sentence[-self.overlap_limit:]

        # Check if current text already starts with the overlap
        if current_text.startswith(last_sentence):
            return current_text

        return last_sentence + " " + current_text


def ensure_directory(directory_path):
    if not os.path.exists(directory_path):
        try:
            os.makedirs(directory_path, mode=0o755)
            logger.info(f"Created directory: {directory_path}")
        except Exception as e:
            logger.error(f"Failed to create directory {directory_path}: {e}")
            return False

    try:
        test_file = os.path.join(directory_path, ".write_test")
        with open(test_file, 'w') as f:
            f.write("test")
        os.remove(test_file)
        return True
    except Exception as e:
        logger.error(f"Directory {directory_path} is not writable: {e}")
        return False


def process_document(file_path, languages=None, segment_size=1000, overlap=True, overlap_limit=200, clean_for_ai=True,
                     segmentation_type="semantic", replace_whitespace=False, remove_urls_emails=False):
    if languages is None:
        languages = ['ch_sim', 'en']

    parser = SemanticDocumentParser(
        language_list=languages,
        segment_size=segment_size,
        overlap=overlap,
        overlap_limit=overlap_limit,
        clean_for_ai=clean_for_ai,
        segmentation_type=segmentation_type,
        replace_whitespace=replace_whitespace,
        remove_urls_emails=remove_urls_emails
    )

    segments = parser.parse(file_path)
    return segments


#################################################
# Dialogue Generation Component
#################################################

class ParallelDialogueGenerator:
    def __init__(self, api_url: str, api_key: str, output_dir: str = "dialogues", model: str = "gpt-4o-mini",
                 max_generation_mode: bool = False, max_questions: int = 10, dialogue_token_limit: int = 1000,
                 language: str = "en-zh"):
        """
        Initialize the dialogue generator, supporting bilingual dialogues.

        Args:
            api_url: Base API URL
            api_key: API authentication key
            output_dir: Directory to save dialogue results
            model: Model to use (default: gpt-4o-mini)
            max_generation_mode: Enable maximum generation mode
            max_questions: Maximum number of questions to generate
            dialogue_token_limit: Maximum token count per dialogue round
            language: Language mode ("en", "zh", or "en-zh" for bilingual)
        """
        # Ensure API URL ends with /v1
        self.api_base = api_url.rstrip('/')
        if not self.api_base.endswith('/v1'):
            if '/v1' not in self.api_base:
                self.api_base += '/v1'

        # Complete chat completion endpoint
        self.api_url = f"{self.api_base}/chat/completions"
        self.api_key = api_key
        self.model = model
        self.output_dir = output_dir
        self.language = language  # Store language preference

        # Parameters for "maximum generation" mode
        self.max_generation_mode = max_generation_mode
        self.max_questions = max_questions
        self.dialogue_token_limit = dialogue_token_limit

        # Create output directory if it doesn't exist
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def _call_api(self, messages: List[Dict[str, str]], max_tokens: int = None, max_retries: int = 3) -> Dict[str, Any]:
        """
        Call the API to send messages.

        Args:
            messages: List of message dictionaries with 'role' and 'content' keys
            max_tokens: Maximum token count for response (overrides default)
            max_retries: Maximum number of API call retries

        Returns:
            API response dictionary
        """
        headers = {
            "Content-Type": "application/json; charset=utf-8",  # Explicitly specify UTF-8 charset
            "Authorization": f"Bearer {self.api_key}"
        }

        # Use provided max_tokens value or class-scope dialogue_token_limit (if in max_generation_mode)
        if max_tokens is None:
            max_tokens = self.dialogue_token_limit if self.max_generation_mode else 1000

        data = {
            "model": self.model,
            "messages": messages,
            "max_tokens": max_tokens,
            "temperature": 0.7
        }

        for attempt in range(max_retries):
            try:
                logger.info(f"Attempt {attempt + 1}: Calling API at {self.api_url}")

                # Manually encode data as JSON with UTF-8 encoding
                json_data = json.dumps(data, ensure_ascii=False).encode('utf-8')

                response = requests.post(
                    self.api_url,
                    headers=headers,
                    data=json_data,  # Use data parameter instead of json
                    timeout=30
                )

                # Print response status and content for debugging
                logger.info(f"Response status: {response.status_code}")
                if response.status_code != 200:
                    logger.error(f"Error response: {response.text[:200]}...")

                response.raise_for_status()
                result = response.json()

                # Validate response format
                if "choices" not in result:
                    raise ValueError(f"Invalid API response format: 'choices' not found in response")

                return result
            except (requests.exceptions.RequestException, json.JSONDecodeError, ValueError) as e:
                if attempt == max_retries - 1:
                    raise Exception(f"API call failed after {max_retries} attempts: {str(e)}")
                logger.error(f"API call attempt {attempt + 1} failed: {str(e)}. Retrying...")
                time.sleep(2 ** attempt)  # Exponential backoff

    def generate_questions(self, chunk: str, num_questions: int = None) -> List[str]:
        """
        Generate questions in the specified language based on the given text chunk.

        Args:
            chunk: Text content to generate questions from
            num_questions: Number of questions to generate (overrides default)

        Returns:
            List of generated questions
        """
        # Determine number of questions to generate
        if num_questions is None:
            num_questions = self.max_questions if self.max_generation_mode else 3

        # Customize prompt based on language setting
        if self.language == "en":
            language_instruction = "Generate questions in English only."
        elif self.language == "zh":
            language_instruction = "请仅用中文生成问题。"
        else:  # en-zh bilingual mode
            language_instruction = "Generate each question in both English and Chinese (同时用英文和中文生成每个问题)."

        prompt = f"""Please generate {num_questions} insightful questions based on the following text.
        Focus on key points, implications, and potential applications of the information.
        {language_instruction}

        {chunk}

        Return only questions in a numbered list form, without any additional text."""

        messages = [
            {"role": "system",
             "content": "You are an assistant who generates insightful questions in both English and Chinese."},
            {"role": "user", "content": prompt}
        ]

        # If in max_generation_mode, allow more tokens for question generation
        max_tokens = 2000 if self.max_generation_mode else 1000
        response = self._call_api(messages, max_tokens=max_tokens)

        # Extract questions from response
        questions_text = response.get("choices", [{}])[0].get("message", {}).get("content", "")

        # Parse questions (supporting both English and Chinese numbering formats)
        questions = []
        for line in questions_text.strip().split('\n'):
            line = line.strip()
            if not line:
                continue

            # Check various numbering formats (1., 1、, -, •, etc.)
            if (line[0].isdigit() or
                    line.startswith('- ') or
                    line.startswith('• ') or
                    line.startswith('* ')):

                # Remove numbering or bullet points
                if '. ' in line and line[0].isdigit():
                    question = line.split('. ', 1)[-1]
                elif '、' in line and line[0].isdigit():  # Chinese style numbering
                    question = line.split('、', 1)[-1]
                elif line.startswith(('- ', '• ', '* ')):
                    question = line[2:]
                else:
                    question = line

                questions.append(question)
            elif questions:  # If this line has no number but we already have questions
                # Might be continuation of previous question
                questions[-1] += " " + line

        return questions[:num_questions]  # Ensure we only return the requested number

    def conduct_dialogue(self, question: str, chunk: str, rounds: int = 3) -> List[Dict[str, str]]:
        """
        Conduct a multi-round dialogue based on the question and context.

        Args:
            question: Initial question
            chunk: Context text block
            rounds: Number of dialogue rounds

        Returns:
            List of messages representing the dialogue
        """
        # Create language-specific system prompt
        if self.language == "en":
            system_content = f"You are having a conversation about the following text in English: {chunk}"
        elif self.language == "zh":
            system_content = f"你正在用中文讨论以下文本内容: {chunk}"
        else:  # en-zh bilingual mode
            system_content = f"You are having a bilingual conversation about the following text. Please reply in the same language as the user's question: {chunk}"

        # Initialize system context and initial question
        messages = [
            {"role": "system", "content": system_content},
            {"role": "user", "content": question}
        ]

        conversation_history = messages.copy()

        # Conduct multi-round dialogue
        for i in range(rounds):
            print(f"Dialogue round {i + 1}/{rounds} for question: {question[:50]}...")

            # Get assistant response
            response = self._call_api(conversation_history)  # Uses class dialogue_token_limit if in max_generation_mode
            assistant_message = {
                "role": "assistant",
                "content": response.get("choices", [{}])[0].get("message", {}).get("content", "")
            }
            conversation_history.append(assistant_message)

            # If this is not the last round, generate follow-up question
            if i < rounds - 1:
                # Adjust follow-up instruction based on language
                if self.language == "en":
                    follow_up_instruction = "Based on the previous exchange, generate a natural follow-up question in English to deepen the conversation."
                elif self.language == "zh":
                    follow_up_instruction = "根据前面的对话，用中文生成一个自然的后续问题，以深入讨论。"
                else:  # en-zh bilingual mode
                    follow_up_instruction = "Based on the previous exchange, generate a natural follow-up question to deepen the conversation. Use the same language (English or Chinese) as the most recent reply."

                follow_up_messages = [
                    {"role": "system", "content": follow_up_instruction},
                    {"role": "user",
                     "content": f"Generate a follow-up question for this response: {assistant_message['content']}"}
                ]
                follow_up_response = self._call_api(follow_up_messages)
                follow_up_question = follow_up_response.get("choices", [{}])[0].get("message", {}).get("content", "")

                # Add follow-up question as user message
                user_follow_up = {"role": "user", "content": follow_up_question}
                conversation_history.append(user_follow_up)

        return conversation_history

    def process_chunk(self, chunk: str, chunk_id: int, rounds: int = 3) -> Dict[str, Any]:
        """
        Process a single text chunk: generate questions and conduct dialogues.

        Args:
            chunk: Text content to process
            chunk_id: Text chunk identifier
            rounds: Number of dialogue rounds

        Returns:
            Dictionary containing the text chunk, questions, and dialogues
        """
        try:
            print(f"Processing text chunk {chunk_id}...")

            # Generate questions based on the text chunk
            questions = self.generate_questions(chunk)
            print(f"Generated {len(questions)} questions for chunk {chunk_id}")

            # Conduct dialogues for each question
            dialogues = {}
            for i, question in enumerate(questions):
                print(f"Conducting dialogue {i + 1}/{len(questions)} for chunk {chunk_id}")
                dialogue = self.conduct_dialogue(question, chunk, rounds)
                dialogues[f"question_{i + 1}"] = {
                    "question": question,
                    "dialogue": dialogue
                }

            # Save results
            result = {
                "chunk_id": chunk_id,
                "chunk_text": chunk,
                "questions": questions,
                "dialogues": dialogues
            }

            self._save_dialogue(result, chunk_id)
            print(f"Successfully completed processing for chunk {chunk_id}")

            return result
        except Exception as e:
            print(f"Error processing chunk {chunk_id}: {str(e)}")
            return {
                "chunk_id": chunk_id,
                "error": str(e),
                "status": "failed"
            }

    def process_chunks_parallel(self, chunks: List[str], rounds: int = 3, max_workers: int = 4) -> List[Dict[str, Any]]:
        """
        Process multiple text chunks in parallel.

        Args:
            chunks: List of text chunks to process
            rounds: Number of dialogue rounds
            max_workers: Maximum number of parallel worker threads

        Returns:
            List of results for each text chunk
        """
        results = []

        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit tasks
            future_to_chunk = {
                executor.submit(self.process_chunk, chunk, i, rounds): (i, chunk)
                for i, chunk in enumerate(chunks)
            }

            # Collect results
            for future in concurrent.futures.as_completed(future_to_chunk):
                chunk_id, chunk = future_to_chunk[future]
                try:
                    result = future.result()
                    results.append(result)
                    if "status" in result and result["status"] == "failed":
                        print(f"Failed processing chunk {chunk_id}: {result.get('error', 'Unknown error')}")
                    else:
                        print(f"Completed processing for chunk {chunk_id}")
                except Exception as e:
                    print(f"Error processing chunk {chunk_id}: {str(e)}")
                    results.append({
                        "chunk_id": chunk_id,
                        "error": str(e),
                        "status": "failed"
                    })

        # Sort results by chunk_id
        results.sort(key=lambda x: x.get("chunk_id", 0))

        return results

    def _save_dialogue(self, dialogue_data: Dict[str, Any], chunk_id: int) -> str:
        """
        Save dialogue data to file.

        Args:
            dialogue_data: Dialogue data to save
            chunk_id: Text chunk identifier

        Returns:
            Path to the saved file
        """
        filename = f"dialogue_chunk_{chunk_id}.json"
        filepath = os.path.join(self.output_dir, filename)

        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(dialogue_data, f, ensure_ascii=False, indent=2)

        return filepath


#################################################
# Integrated Workflow
#################################################

def process_document_and_generate_dialogues(
        file_path,
        api_url,
        api_key,
        model="gpt-4o-mini",
        dialogue_rounds=3,
        max_workers=4,
        output_dir="results",
        language="en-zh",
        max_generation_mode=False,
        max_questions=5,
        dialogue_token_limit=1000,
        segment_size=1000,
        overlap=True,
        overlap_limit=200,
        clean_for_ai=True,
        segmentation_type="semantic",
        replace_whitespace=False,
        remove_urls_emails=False,
        languages=None
):
    """
    Integrated workflow: Process document into segments and generate dialogues.

    Args:
        file_path: Path to document file
        api_url: API base URL for dialogue generation
        api_key: API key for dialogue generation
        model: Model to use for dialogue generation
        dialogue_rounds: Number of dialogue rounds
        max_workers: Maximum parallel workers
        output_dir: Directory to save results
        language: Language for dialogues (en, zh, en-zh)
        max_generation_mode: Enable max generation mode
        max_questions: Maximum questions per segment
        dialogue_token_limit: Token limit per dialogue round
        segment_size: Maximum size of each segment
        overlap: Enable segment overlap
        overlap_limit: Maximum overlap between segments
        clean_for_ai: Clean text for AI training
        segmentation_type: Type of segmentation (semantic or structural)
        replace_whitespace: Replace consecutive whitespace
        remove_urls_emails: Remove URLs and emails
        languages: Languages for document processing

    Returns:
        Dictionary with processing results
    """
    if languages is None:
        languages = ['ch_sim', 'en']

    # Ensure output directory exists
    output_dir = os.path.abspath(output_dir)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    dialogue_dir = os.path.join(output_dir, "dialogues")
    if not os.path.exists(dialogue_dir):
        os.makedirs(dialogue_dir)

    # Step 1: Parse document into segments
    logger.info(f"Processing document: {file_path}")
    try:
        segments = process_document(
            file_path,
            languages=languages,
            segment_size=segment_size,
            overlap=overlap,
            overlap_limit=overlap_limit,
            clean_for_ai=clean_for_ai,
            segmentation_type=segmentation_type,
            replace_whitespace=replace_whitespace,
            remove_urls_emails=remove_urls_emails
        )

        logger.info(f"Document processed into {len(segments)} segments")

        # Save segments to file
        segments_file = os.path.join(output_dir, "segments.json")
        with open(segments_file, 'w', encoding='utf-8') as f:
            json.dump(segments, f, ensure_ascii=False, indent=2)

        # Step 2: Generate dialogues from segments
        chunks = [segment['content'] for segment in segments]

        generator = ParallelDialogueGenerator(
            api_url=api_url,
            api_key=api_key,
            output_dir=dialogue_dir,
            model=model,
            max_generation_mode=max_generation_mode,
            max_questions=max_questions,
            dialogue_token_limit=dialogue_token_limit,
            language=language
        )

        logger.info(f"Generating dialogues for {len(chunks)} segments")
        dialogue_results = generator.process_chunks_parallel(
            chunks=chunks,
            rounds=dialogue_rounds,
            max_workers=max_workers
        )

        # Save combined results
        combined_results = {
            "document": os.path.basename(file_path),
            "segments": segments,
            "dialogues": dialogue_results
        }

        results_file = os.path.join(output_dir, "combined_results.json")
        with open(results_file, 'w', encoding='utf-8') as f:
            json.dump(combined_results, f, ensure_ascii=False, indent=2)

        return {
            "status": "success",
            "segments_count": len(segments),
            "dialogues_count": len(dialogue_results),
            "results_file": results_file
        }

    except Exception as e:
        logger.error(f"Error in integrated workflow: {str(e)}")
        return {
            "status": "error",
            "error": str(e)
        }


#################################################
# Flask Web Application
#################################################

app = Flask(__name__)
CORS(app)

# Configure maximum content length (50MB by default)
MAX_CONTENT_LENGTH = int(os.environ.get('MAX_CONTENT_LENGTH', 50 * 1024 * 1024))
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# Set up upload folder
UPLOAD_FOLDER = os.environ.get('UPLOAD_FOLDER')
upload_paths = [
    UPLOAD_FOLDER,
    '/tmp/uploads',
    '/var/tmp/uploads',
    os.path.join(os.getcwd(), 'uploads'),
    tempfile.gettempdir()
]

for path in upload_paths:
    if path and ensure_directory(path):
        UPLOAD_FOLDER = path
        logger.info(f"Using upload folder: {UPLOAD_FOLDER}")
        break
else:
    UPLOAD_FOLDER = tempfile.mkdtemp()
    logger.warning(f"Falling back to temporary directory: {UPLOAD_FOLDER}")

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Set up output folder for results
RESULTS_FOLDER = os.environ.get('RESULTS_FOLDER', 'results')
if not os.path.exists(RESULTS_FOLDER):
    os.makedirs(RESULTS_FOLDER)
app.config['RESULTS_FOLDER'] = RESULTS_FOLDER


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/api/process-document', methods=['POST'])
def api_process_document():
    """API endpoint to process a document into segments"""
    logger.info("Received document processing request")

    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    # Process form parameters
    try:
        languages = request.form.get('languages', 'ch_sim,en').split(',')
        segment_size = int(request.form.get('segment_size', 1000))
        overlap = request.form.get('overlap', 'true').lower() == 'true'
        overlap_limit = int(request.form.get('overlap_limit', 200))
        clean_for_ai = request.form.get('clean_for_ai', 'true').lower() == 'true'
        segmentation_type = request.form.get('segmentation_type', 'semantic')
        replace_whitespace = request.form.get('replace_whitespace', 'false').lower() == 'true'
        remove_urls_emails = request.form.get('remove_urls_emails', 'false').lower() == 'true'
    except Exception as e:
        return jsonify({"error": f"Invalid request parameters: {str(e)}"}), 400

    # Save file
    try:
        filename = str(uuid.uuid4()) + '_' + secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)

        # Process document
        segments = process_document(
            file_path,
            languages=languages,
            segment_size=segment_size,
            overlap=overlap,
            overlap_limit=overlap_limit,
            clean_for_ai=clean_for_ai,
            segmentation_type=segmentation_type,
            replace_whitespace=replace_whitespace,
            remove_urls_emails=remove_urls_emails
        )

        # Clean up file
        try:
            os.remove(file_path)
        except Exception as e:
            logger.warning(f"Failed to remove temporary file {file_path}: {e}")

        return jsonify({
            "status": "success",
            "segments_count": len(segments),
            "segments": segments
        })

    except Exception as e:
        logger.error(f"Error processing document: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/generate-dialogues', methods=['POST'])
def api_generate_dialogues():
    """API endpoint to generate dialogues from segments"""
    logger.info("Received dialogue generation request")

    if not request.json or 'segments' not in request.json:
        return jsonify({"error": "No segments provided"}), 400

    # Process request parameters
    try:
        segments = request.json['segments']
        api_url = request.json.get('api_url')
        api_key = request.json.get('api_key')

        if not api_url or not api_key:
            return jsonify({"error": "API URL and API key are required"}), 400

        model = request.json.get('model', 'gpt-4o-mini')
        dialogue_rounds = int(request.json.get('dialogue_rounds', 3))
        max_workers = int(request.json.get('max_workers', 4))
        language = request.json.get('language', 'en-zh')
        max_generation_mode = request.json.get('max_generation_mode', False)
        max_questions = int(request.json.get('max_questions', 5))
        dialogue_token_limit = int(request.json.get('dialogue_token_limit', 1000))
    except Exception as e:
        return jsonify({"error": f"Invalid request parameters: {str(e)}"}), 400

    # Create unique directory for results
    result_dir = os.path.join(app.config['RESULTS_FOLDER'], str(uuid.uuid4()))
    if not os.path.exists(result_dir):
        os.makedirs(result_dir)

    # Extract text content from segments
    chunks = [segment['content'] for segment in segments]

    try:
        # Initialize dialogue generator
        generator = ParallelDialogueGenerator(
            api_url=api_url,
            api_key=api_key,
            output_dir=result_dir,
            model=model,
            max_generation_mode=max_generation_mode,
            max_questions=max_questions,
            dialogue_token_limit=dialogue_token_limit,
            language=language
        )

        # Process chunks
        dialogue_results = generator.process_chunks_parallel(
            chunks=chunks,
            rounds=dialogue_rounds,
            max_workers=max_workers
        )

        # Save combined results
        combined_results = {
            "segments": segments,
            "dialogues": dialogue_results
        }

        results_file = os.path.join(result_dir, "results.json")
        with open(results_file, 'w', encoding='utf-8') as f:
            json.dump(combined_results, f, ensure_ascii=False, indent=2)

        return jsonify({
            "status": "success",
            "dialogues_count": len(dialogue_results),
            "dialogues": dialogue_results,
            "results_file": results_file
        })

    except Exception as e:
        logger.error(f"Error generating dialogues: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/integrated-workflow', methods=['POST'])
def api_integrated_workflow():
    """API endpoint for integrated document processing and dialogue generation"""
    logger.info("Received integrated workflow request")

    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    # Check for required parameters
    api_url = request.form.get('api_url')
    api_key = request.form.get('api_key')

    if not api_url or not api_key:
        return jsonify({"error": "API URL and API key are required"}), 400

    # Process other parameters
    try:
        # Document parsing parameters
        languages = request.form.get('languages', 'ch_sim,en').split(',')
        segment_size = int(request.form.get('segment_size', 1000))
        overlap = request.form.get('overlap', 'true').lower() == 'true'
        overlap_limit = int(request.form.get('overlap_limit', 200))
        clean_for_ai = request.form.get('clean_for_ai', 'true').lower() == 'true'
        segmentation_type = request.form.get('segmentation_type', 'semantic')
        replace_whitespace = request.form.get('replace_whitespace', 'false').lower() == 'true'
        remove_urls_emails = request.form.get('remove_urls_emails', 'false').lower() == 'true'

        # Dialogue generation parameters
        model = request.form.get('model', 'gpt-4o-mini')
        dialogue_rounds = int(request.form.get('dialogue_rounds', 3))
        max_workers = int(request.form.get('max_workers', 4))
        language = request.form.get('language', 'en-zh')
        max_generation_mode = request.form.get('max_generation_mode', 'false').lower() == 'true'
        max_questions = int(request.form.get('max_questions', 5))
        dialogue_token_limit = int(request.form.get('dialogue_token_limit', 1000))
    except Exception as e:
        return jsonify({"error": f"Invalid request parameters: {str(e)}"}), 400

    # Create unique directory for results
    result_dir = os.path.join(app.config['RESULTS_FOLDER'], str(uuid.uuid4()))
    if not os.path.exists(result_dir):
        os.makedirs(result_dir)

    # Save file
    try:
        filename = str(uuid.uuid4()) + '_' + secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)

        # Run integrated workflow
        result = process_document_and_generate_dialogues(
            file_path=file_path,
            api_url=api_url,
            api_key=api_key,
            model=model,
            dialogue_rounds=dialogue_rounds,
            max_workers=max_workers,
            output_dir=result_dir,
            language=language,
            max_generation_mode=max_generation_mode,
            max_questions=max_questions,
            dialogue_token_limit=dialogue_token_limit,
            segment_size=segment_size,
            overlap=overlap,
            overlap_limit=overlap_limit,
            clean_for_ai=clean_for_ai,
            segmentation_type=segmentation_type,
            replace_whitespace=replace_whitespace,
            remove_urls_emails=remove_urls_emails,
            languages=languages
        )

        # Clean up temporary file
        try:
            os.remove(file_path)
        except Exception as e:
            logger.warning(f"Failed to remove temporary file {file_path}: {e}")

        # Return result
        if result["status"] == "success":
            # Read the combined results file
            with open(result["results_file"], 'r', encoding='utf-8') as f:
                combined_results = json.load(f)

            return jsonify({
                "status": "success",
                "segments_count": result["segments_count"],
                "dialogues_count": result["dialogues_count"],
                "results": combined_results,
                "results_file": result["results_file"]
            })
        else:
            return jsonify({"error": result["error"]}), 500

    except Exception as e:
        logger.error(f"Error in integrated workflow: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/download-results', methods=['POST'])
def api_download_results():
    """API endpoint to download results as JSON file"""
    logger.info("Received download results request")

    if not request.json:
        return jsonify({"error": "No JSON data provided"}), 400

    try:
        fd, temp_file_path = tempfile.mkstemp(suffix='.json')
        with os.fdopen(fd, 'w', encoding='utf-8') as f:
            json.dump(request.json, f, ensure_ascii=False, indent=2)

        return send_file(
            temp_file_path,
            as_attachment=True,
            download_name="document_dialogue_results.json",
            mimetype='application/json'
        )
    except Exception as e:
        logger.error(f"Error creating download file: {e}")
        return jsonify({"error": f"Error creating download file: {str(e)}"}), 500


@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    status = {
        "status": "healthy",
        "version": "1.0.0",
        "upload_folder": app.config['UPLOAD_FOLDER'],
        "results_folder": app.config['RESULTS_FOLDER'],
        "upload_folder_writable": ensure_directory(app.config['UPLOAD_FOLDER']),
        "results_folder_writable": ensure_directory(app.config['RESULTS_FOLDER']),
        "temp_dir_writable": ensure_directory(tempfile.gettempdir()),
        "optional_dependencies": OPTIONAL_DEPS_AVAILABLE
    }
    return jsonify(status)


#################################################
# Main Function
#################################################

def main():
    parser = argparse.ArgumentParser(description='Document Analysis and Dialogue Generation Tool')
    subparsers = parser.add_subparsers(dest='command', help='Command')

    # Server command
    parser_server = subparsers.add_parser('server', help='Start web server')
    parser_server.add_argument('--host', default='0.0.0.0', help='Server host address')
    parser_server.add_argument('--port', type=int, default=5000, help='Server port')
    parser_server.add_argument('--debug', action='store_true', help='Enable debug mode')

    # Document processing command
    parser_document = subparsers.add_parser('document', help='Process document into segments')
    parser_document.add_argument('file_path', help='Path to document file')
    parser_document.add_argument('--output', help='Output file path for results (JSON format)')
    parser_document.add_argument('--languages', nargs='+', default=['ch_sim', 'en'],
                                 help='OCR languages (e.g., ch_sim en ja)')
    parser_document.add_argument('--segment-size', type=int, default=1000,
                                 help='Maximum characters per segment')
    parser_document.add_argument('--no-overlap', action='store_true',
                                 help='Disable segment overlap (enabled by default)')
    parser_document.add_argument('--overlap-limit', type=int, default=200,
                                 help='Maximum character overlap')
    parser_document.add_argument('--no-clean-ai', action='store_true',
                                 help='Disable AI training data cleaning (enabled by default)')
    parser_document.add_argument('--segmentation-type', choices=['semantic', 'structural'], default='semantic',
                                 help='Segmentation method: semantic or structural')
    parser_document.add_argument('--replace-whitespace', action='store_true',
                                 help='Replace consecutive whitespace with single space')
    parser_document.add_argument('--remove-urls-emails', action='store_true',
                                 help='Remove all URLs and email addresses')

    # Dialogue generation command
    parser_dialogue = subparsers.add_parser('dialogue', help='Generate dialogues from text chunks')
    parser_dialogue.add_argument('input_file', help='JSON file containing text chunks')
    parser_dialogue.add_argument('--api-url', required=True, help='Base API URL')
    parser_dialogue.add_argument('--api-key', required=True, help='API authentication key')
    parser_dialogue.add_argument('--output-dir', default='dialogues', help='Directory to save dialogue results')
    parser_dialogue.add_argument('--model', default='gpt-4o-mini', help='Model to use')
    parser_dialogue.add_argument('--rounds', type=int, default=3, help='Number of dialogue rounds')
    parser_dialogue.add_argument('--max-workers', type=int, default=4, help='Maximum parallel worker threads')
    parser_dialogue.add_argument('--language', default='en-zh', choices=['en', 'zh', 'en-zh'],
                                 help='Language mode: English (en), Chinese (zh), or bilingual (en-zh)')
    parser_dialogue.add_argument('--max-generation-mode', action='store_true',
                                 help='Enable maximum generation mode')
    parser_dialogue.add_argument('--max-questions', type=int, default=10,
                                 help='Maximum questions to generate')
    parser_dialogue.add_argument('--token-limit', type=int, default=1000,
                                 help='Maximum tokens per dialogue round')

    # Integrated workflow command
    parser_integrated = subparsers.add_parser('integrated', help='Run integrated document and dialogue workflow')
    parser_integrated.add_argument('file_path', help='Path to document file')
    parser_integrated.add_argument('--api-url', required=True, help='Base API URL')
    parser_integrated.add_argument('--api-key', required=True, help='API authentication key')
    parser_integrated.add_argument('--output-dir', default='results', help='Directory to save results')
    parser_integrated.add_argument('--model', default='gpt-4o-mini', help='Model to use')
    parser_integrated.add_argument('--dialogue-rounds', type=int, default=3, help='Number of dialogue rounds')
    parser_integrated.add_argument('--max-workers', type=int, default=4, help='Maximum parallel worker threads')
    parser_integrated.add_argument('--language', default='en-zh', choices=['en', 'zh', 'en-zh'],
                                   help='Language mode: English (en), Chinese (zh), or bilingual (en-zh)')
    parser_integrated.add_argument('--max-generation-mode', action='store_true',
                                   help='Enable maximum generation mode')
    parser_integrated.add_argument('--max-questions', type=int, default=5,
                                   help='Maximum questions to generate per segment')
    parser_integrated.add_argument('--token-limit', type=int, default=1000,
                                   help='Maximum tokens per dialogue round')
    parser_integrated.add_argument('--languages', nargs='+', default=['ch_sim', 'en'],
                                   help='OCR languages (e.g., ch_sim en ja)')
    parser_integrated.add_argument('--segment-size', type=int, default=1000,
                                   help='Maximum characters per segment')
    parser_integrated.add_argument('--no-overlap', action='store_true',
                                   help='Disable segment overlap (enabled by default)')
    parser_integrated.add_argument('--overlap-limit', type=int, default=200,
                                   help='Maximum character overlap')
    parser_integrated.add_argument('--no-clean-ai', action='store_true',
                                   help='Disable AI training data cleaning (enabled by default)')
    parser_integrated.add_argument('--segmentation-type', choices=['semantic', 'structural'], default='semantic',
                                   help='Segmentation method: semantic or structural')

    args = parser.parse_args()

    if args.command == 'server':
        # Create HTML template
        template_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'templates')
        if not os.path.exists(template_dir):
            os.makedirs(template_dir)

        with open(os.path.join(template_dir, 'index.html'), 'w', encoding='utf-8') as f:
            f.write('''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Document Analysis & Dialogue Generation</title>
    <style>
        body {
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            line-height: 1.6;
            color: #333;
            max-width: 1200px;
            margin: 0 auto;
            padding: 20px;
        }
        .card {
            background: #fff;
            border-radius: 8px;
            box-shadow: 0 2px 10px rgba(0, 0, 0, 0.1);
            padding: 20px;
            margin-bottom: 20px;
        }
        h1, h2, h3 {
            color: #2c3e50;
        }
        .tabs {
            display: flex;
            margin-bottom: 20px;
            border-bottom: 1px solid #ddd;
        }
        .tab {
            padding: 10px 20px;
            cursor: pointer;
            border-bottom: 3px solid transparent;
        }
        .tab.active {
            border-bottom: 3px solid #3498db;
            font-weight: bold;
            color: #3498db;
        }
        .tab-content {
            display: none;
        }
        .tab-content.active {
            display: block;
        }
        .form-group {
            margin-bottom: 15px;
        }
        label {
            display: block;
            margin-bottom: 5px;
            font-weight: 500;
        }
        input[type="text"], input[type="number"], input[type="password"], select, textarea {
            width: 100%;
            padding: 8px;
            border: 1px solid #ddd;
            border-radius: 4px;
            box-sizing: border-box;
        }
        input[type="checkbox"] {
            margin-right: 8px;
        }
        button {
            background: #3498db;
            color: white;
            border: none;
            padding: 10px 20px;
            border-radius: 4px;
            cursor: pointer;
            font-size: 16px;
        }
        button:hover {
            background: #2980b9;
        }
        button.secondary {
            background: #95a5a6;
        }
        button.secondary:hover {
            background: #7f8c8d;
        }
        .loading {
            display: none;
            margin: 20px 0;
            text-align: center;
        }
        .spinner {
            border: 4px solid rgba(0, 0, 0, 0.1);
            width: 36px;
            height: 36px;
            border-radius: 50%;
            border-left-color: #3498db;
            animation: spin 1s linear infinite;
            display: inline-block;
            vertical-align: middle;
            margin-right: 10px;
        }
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        .panel {
            border: 1px solid #ddd;
            border-radius: 4px;
            margin-bottom: 15px;
        }
        .panel-header {
            background: #f5f5f5;
            padding: 10px 15px;
            border-bottom: 1px solid #ddd;
            font-weight: 500;
            cursor: pointer;
            display: flex;
            justify-content: space-between;
        }
        .panel-body {
            padding: 15px;
            display: none;
        }
        .panel-body.active {
            display: block;
        }
        .badge {
            background: #3498db;
            color: white;
            border-radius: 10px;
            padding: 2px 8px;
            font-size: 12px;
        }
        .toggle-btn {
            background: none;
            border: none;
            padding: 0;
            font-size: 20px;
            cursor: pointer;
            color: #7f8c8d;
        }
        .grid {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 20px;
        }
        .hidden {
            display: none;
        }
        pre {
            background: #f8f8f8;
            padding: 10px;
            border-radius: 4px;
            overflow-x: auto;
            font-size: 14px;
        }
        .dialogue {
            margin-bottom: 20px;
        }
        .message {
            padding: 10px 15px;
            border-radius: 10px;
            margin-bottom: 10px;
            max-width: 80%;
        }
        .user {
            background: #f1f0f0;
            margin-left: auto;
        }
        .assistant {
            background: #e3f2fd;
            margin-right: auto;
        }
        .results {
            margin-top: 20px;
            display: none;
        }
        .error {
            color: #e74c3c;
            font-weight: bold;
        }
        .section-header {
            display: flex;
            justify-content: space-between;
            align-items: center;
        }
        .two-columns {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 20px;
        }
        @media (max-width: 768px) {
            .grid, .two-columns {
                grid-template-columns: 1fr;
            }
        }
    </style>
</head>
<body>
    <div class="card">
        <h1>Document Analysis & Dialogue Generation</h1>
        <p>Upload a document to analyze it, generate questions, and conduct AI dialogues.</p>

        <div class="tabs">
            <div class="tab active" data-tab="integrated">Integrated Workflow</div>
            <div class="tab" data-tab="parsing">Document Parsing</div>
            <div class="tab" data-tab="dialogue">Dialogue Generation</div>
        </div>

        <!-- Integrated Workflow Tab -->
        <div class="tab-content active" id="integrated-tab">
            <form id="integrated-form">
                <div class="grid">
                    <div>
                        <h3>Document Processing</h3>
                        <div class="form-group">
                            <label for="integrated-file">Select Document:</label>
                            <input type="file" id="integrated-file" name="file" required>
                        </div>

                        <div class="panel">
                            <div class="panel-header">
                                Document Processing Settings
                                <button type="button" class="toggle-btn">+</button>
                            </div>
                            <div class="panel-body">
                                <div class="form-group">
                                    <label for="integrated-segmentation-type">Segmentation Method:</label>
                                    <select id="integrated-segmentation-type" name="segmentation_type">
                                        <option value="semantic">Semantic (by meaning)</option>
                                        <option value="structural">Structural (by document structure)</option>
                                    </select>
                                </div>

                                <div class="form-group">
                                    <label for="integrated-languages">OCR Languages (comma separated):</label>
                                    <input type="text" id="integrated-languages" name="languages" value="ch_sim,en">
                                </div>

                                <div class="form-group">
                                    <label for="integrated-segment-size">Maximum Segment Size (characters):</label>
                                    <input type="number" id="integrated-segment-size" name="segment_size" value="1000">
                                </div>

                                <div class="form-group">
                                    <label>
                                        <input type="checkbox" id="integrated-overlap" name="overlap" checked>
                                        Enable Semantic Overlap
                                    </label>
                                </div>

                                <div class="form-group">
                                    <label for="integrated-overlap-limit">Overlap Limit (characters):</label>
                                    <input type="number" id="integrated-overlap-limit" name="overlap_limit" value="200">
                                </div>

                                <div class="form-group">
                                    <label>
                                        <input type="checkbox" id="integrated-clean-ai" name="clean_for_ai" checked>
                                        Clean Text for AI Processing
                                    </label>
                                </div>
                            </div>
                        </div>
                    </div>

                    <div>
                        <h3>Dialogue Generation</h3>
                        <div class="form-group">
                            <label for="integrated-api-url">API URL:</label>
                            <input type="text" id="integrated-api-url" name="api_url" placeholder="https://api.example.com/v1" required>
                        </div>

                        <div class="form-group">
                            <label for="integrated-api-key">API Key:</label>
                            <input type="password" id="integrated-api-key" name="api_key" required>
                        </div>

                        <div class="panel">
                            <div class="panel-header">
                                Dialogue Settings
                                <button type="button" class="toggle-btn">+</button>
                            </div>
                            <div class="panel-body">
                                <div class="form-group">
                                    <label for="integrated-model">Model:</label>
                                    <input type="text" id="integrated-model" name="model" value="gpt-4o-mini">
                                </div>

                                <div class="form-group">
                                    <label for="integrated-language">Dialogue Language:</label>
                                    <select id="integrated-language" name="language">
                                        <option value="en-zh">Bilingual (English & Chinese)</option>
                                        <option value="en">English Only</option>
                                        <option value="zh">Chinese Only</option>
                                    </select>
                                </div>

                                <div class="form-group">
                                    <label for="integrated-dialogue-rounds">Dialogue Rounds:</label>
                                    <input type="number" id="integrated-dialogue-rounds" name="dialogue_rounds" value="3">
                                </div>

                                <div class="form-group">
                                    <label>
                                        <input type="checkbox" id="integrated-max-generation" name="max_generation_mode">
                                        Maximum Generation Mode
                                    </label>
                                </div>

                                <div class="form-group">
                                    <label for="integrated-max-questions">Maximum Questions per Segment:</label>
                                    <input type="number" id="integrated-max-questions" name="max_questions" value="5">
                                </div>

                                <div class="form-group">
                                    <label for="integrated-token-limit">Token Limit per Dialogue Round:</label>
                                    <input type="number" id="integrated-token-limit" name="dialogue_token_limit" value="1000">
                                </div>

                                <div class="form-group">
                                    <label for="integrated-max-workers">Maximum Parallel Workers:</label>
                                    <input type="number" id="integrated-max-workers" name="max_workers" value="4">
                                </div>
                            </div>
                        </div>
                    </div>
                </div>

                <button type="submit">Process Document & Generate Dialogues</button>
            </form>

            <div id="integrated-loading" class="loading">
                <div class="spinner"></div>
                <span>Processing document and generating dialogues... This may take several minutes.</span>
            </div>

            <div id="integrated-results" class="results">
                <div class="section-header">
                    <h2>Results</h2>
                    <button id="integrated-download-btn" class="secondary">Download JSON</button>
                </div>

                <div class="card">
                    <h3>Document Segments <span id="integrated-segments-count" class="badge">0</span></h3>
                    <div id="integrated-segments"></div>
                </div>

                <div class="card">
                    <h3>Dialogues <span id="integrated-dialogues-count" class="badge">0</span></h3>
                    <div id="integrated-dialogues"></div>
                </div>
            </div>
        </div>

        <!-- Document Parsing Tab -->
        <div class="tab-content" id="parsing-tab">
            <form id="parsing-form">
                <div class="form-group">
                    <label for="parsing-file">Select Document:</label>
                    <input type="file" id="parsing-file" name="file" required>
                </div>

                <div class="two-columns">
                    <div class="form-group">
                        <label for="parsing-segmentation-type">Segmentation Method:</label>
                        <select id="parsing-segmentation-type" name="segmentation_type">
                            <option value="semantic">Semantic (by meaning)</option>
                            <option value="structural">Structural (by document structure)</option>
                        </select>
                    </div>

                    <div class="form-group">
                        <label for="parsing-languages">OCR Languages (comma separated):</label>
                        <input type="text" id="parsing-languages" name="languages" value="ch_sim,en">
                    </div>
                </div>

                <div class="two-columns">
                    <div class="form-group">
                        <label for="parsing-segment-size">Maximum Segment Size (characters):</label>
                        <input type="number" id="parsing-segment-size" name="segment_size" value="1000">
                    </div>

                    <div class="form-group">
                        <label for="parsing-overlap-limit">Overlap Limit (characters):</label>
                        <input type="number" id="parsing-overlap-limit" name="overlap_limit" value="200">
                    </div>
                </div>

                <div class="two-columns">
                    <div class="form-group">
                        <label>
                            <input type="checkbox" id="parsing-overlap" name="overlap" checked>
                            Enable Semantic Overlap
                        </label>
                    </div>

                    <div class="form-group">
                        <label>
                            <input type="checkbox" id="parsing-clean-ai" name="clean_for_ai" checked>
                            Clean Text for AI Processing
                        </label>
                    </div>
                </div>

                <div class="two-columns">
                    <div class="form-group">
                        <label>
                            <input type="checkbox" id="parsing-replace-whitespace" name="replace_whitespace">
                            Replace Consecutive Whitespace
                        </label>
                    </div>

                    <div class="form-group">
                        <label>
                            <input type="checkbox" id="parsing-remove-urls-emails" name="remove_urls_emails">
                            Remove URLs and Emails
                        </label>
                    </div>
                </div>

                <button type="submit">Process Document</button>
            </form>

            <div id="parsing-loading" class="loading">
                <div class="spinner"></div>
                <span>Processing document... This may take a moment.</span>
            </div>

            <div id="parsing-results" class="results">
                <div class="section-header">
                    <h2>Segments <span id="parsing-segments-count" class="badge">0</span></h2>
                    <button id="parsing-download-btn" class="secondary">Download JSON</button>
                </div>
                <div id="parsing-segments"></div>
            </div>
        </div>

        <!-- Dialogue Generation Tab -->
        <div class="tab-content" id="dialogue-tab">
            <form id="dialogue-form">
                <h3>Dialogue Settings</h3>

                <div class="form-group">
                    <label for="dialogue-segments">Text Segments (JSON array):</label>
                    <textarea id="dialogue-segments" name="segments" rows="6" placeholder='[{"content": "Your text here", "type": "semantic_unit"}]' required></textarea>
                </div>

                <div class="form-group">
                    <label for="dialogue-api-url">API URL:</label>
                    <input type="text" id="dialogue-api-url" name="api_url" placeholder="https://api.example.com/v1" required>
                </div>

                <div class="form-group">
                    <label for="dialogue-api-key">API Key:</label>
                    <input type="password" id="dialogue-api-key" name="api_key" required>
                </div>

                <div class="two-columns">
                    <div class="form-group">
                        <label for="dialogue-model">Model:</label>
                        <input type="text" id="dialogue-model" name="model" value="gpt-4o-mini">
                    </div>

                    <div class="form-group">
                        <label for="dialogue-language">Dialogue Language:</label>
                        <select id="dialogue-language" name="language">
                            <option value="en-zh">Bilingual (English & Chinese)</option>
                            <option value="en">English Only</option>
                            <option value="zh">Chinese Only</option>
                        </select>
                    </div>
                </div>

                <div class="two-columns">
                    <div class="form-group">
                        <label for="dialogue-rounds">Dialogue Rounds:</label>
                        <input type="number" id="dialogue-rounds" name="dialogue_rounds" value="3">
                    </div>

                    <div class="form-group">
                        <label for="dialogue-max-workers">Maximum Parallel Workers:</label>
                        <input type="number" id="dialogue-max-workers" name="max_workers" value="4">
                    </div>
                </div>

                <div class="two-columns">
                    <div class="form-group">
                        <label>
                            <input type="checkbox" id="dialogue-max-generation" name="max_generation_mode">
                            Maximum Generation Mode
                        </label>
                    </div>

                    <div class="form-group">
                        <label for="dialogue-max-questions">Maximum Questions per Segment:</label>
                        <input type="number" id="dialogue-max-questions" name="max_questions" value="5">
                    </div>
                </div>

                <div class="form-group">
                    <label for="dialogue-token-limit">Token Limit per Dialogue Round:</label>
                    <input type="number" id="dialogue-token-limit" name="dialogue_token_limit" value="1000">
                </div>

                <button type="submit">Generate Dialogues</button>
            </form>

            <div id="dialogue-loading" class="loading">
                <div class="spinner"></div>
                <span>Generating dialogues... This may take several minutes.</span>
            </div>

            <div id="dialogue-results" class="results">
                <div class="section-header">
                    <h2>Dialogues <span id="dialogue-dialogues-count" class="badge">0</span></h2>
                    <button id="dialogue-download-btn" class="secondary">Download JSON</button>
                </div>

                <div id="dialogue-content"></div>
            </div>
        </div>
    </div>

    <script>
        // Tab navigation
        document.querySelectorAll('.tab').forEach(tab => {
            tab.addEventListener('click', () => {
                // Deactivate all tabs and tab contents
                document.querySelectorAll('.tab').forEach(t => t.classList.remove('active'));
                document.querySelectorAll('.tab-content').forEach(c => c.classList.remove('active'));

                // Activate clicked tab and corresponding content
                tab.classList.add('active');
                document.getElementById(tab.dataset.tab + '-tab').classList.add('active');
            });
        });

        // Toggle panel sections
        document.querySelectorAll('.panel-header').forEach(header => {
            header.addEventListener('click', () => {
                const body = header.nextElementSibling;
                body.classList.toggle('active');
                const btn = header.querySelector('.toggle-btn');
                btn.textContent = body.classList.contains('active') ? '-' : '+';
            });
        });

        // Integrated workflow form submission
        document.getElementById('integrated-form').addEventListener('submit', async function(e) {
            e.preventDefault();

            const formData = new FormData(this);
            formData.set('overlap', document.getElementById('integrated-overlap').checked);
            formData.set('clean_for_ai', document.getElementById('integrated-clean-ai').checked);
            formData.set('max_generation_mode', document.getElementById('integrated-max-generation').checked);

            document.getElementById('integrated-loading').style.display = 'block';
            document.getElementById('integrated-results').style.display = 'none';

            try {
                const response = await fetch('/api/integrated-workflow', {
                    method: 'POST',
                    body: formData
                });

                const data = await response.json();

                if (!response.ok) {
                    throw new Error(data.error || 'An error occurred');
                }

                // Store result for download
                window.integratedResults = data.results;

                // Show results
                document.getElementById('integrated-segments-count').textContent = data.segments_count;
                document.getElementById('integrated-dialogues-count').textContent = data.dialogues_count;

                // Render segments
                const segmentsContainer = document.getElementById('integrated-segments');
                segmentsContainer.innerHTML = '';

                data.results.segments.forEach((segment, index) => {
                    const segmentDiv = document.createElement('div');
                    segmentDiv.className = 'panel';

                    const header = document.createElement('div');
                    header.className = 'panel-header';
                    header.innerHTML = `Segment ${index + 1} <button type="button" class="toggle-btn">+</button>`;

                    const body = document.createElement('div');
                    body.className = 'panel-body';
                    body.innerHTML = `<pre>${segment.content}</pre>`;

                    segmentDiv.appendChild(header);
                    segmentDiv.appendChild(body);
                    segmentsContainer.appendChild(segmentDiv);

                    header.addEventListener('click', () => {
                        body.classList.toggle('active');
                        const btn = header.querySelector('.toggle-btn');
                        btn.textContent = body.classList.contains('active') ? '-' : '+';
                    });
                });

                // Render dialogues
                const dialoguesContainer = document.getElementById('integrated-dialogues');
                dialoguesContainer.innerHTML = '';

                data.results.dialogues.forEach((result, index) => {
                    const dialogueDiv = document.createElement('div');
                    dialogueDiv.className = 'panel';

                    const header = document.createElement('div');
                    header.className = 'panel-header';
                    header.innerHTML = `Segment ${index + 1} Dialogues <span class="badge">${Object.keys(result.dialogues).length}</span> <button type="button" class="toggle-btn">+</button>`;

                    const body = document.createElement('div');
                    body.className = 'panel-body';

                    // For each question in this chunk
                    Object.entries(result.dialogues).forEach(([qKey, qData]) => {
                        const questionDiv = document.createElement('div');
                        questionDiv.className = 'dialogue';

                        const questionHeader = document.createElement('h4');
                        questionHeader.textContent = qData.question;
                        questionDiv.appendChild(questionHeader);

                        // Display conversation
                        qData.dialogue.forEach(msg => {
                            if (msg.role === 'system') return; // Skip system messages

                            const messageDiv = document.createElement('div');
                            messageDiv.className = `message ${msg.role}`;
                            messageDiv.textContent = msg.content;
                            questionDiv.appendChild(messageDiv);
                        });

                        body.appendChild(questionDiv);
                    });

                    dialogueDiv.appendChild(header);
                    dialogueDiv.appendChild(body);
                    dialoguesContainer.appendChild(dialogueDiv);

                    header.addEventListener('click', () => {
                        body.classList.toggle('active');
                        const btn = header.querySelector('.toggle-btn');
                        btn.textContent = body.classList.contains('active') ? '-' : '+';
                    });
                });

                document.getElementById('integrated-results').style.display = 'block';

            } catch (error) {
                console.error('Error:', error);
                const resultsDiv = document.getElementById('integrated-results');
                resultsDiv.style.display = 'block';
                resultsDiv.innerHTML = `<div class="error">Error: ${error.message}</div>`;
            } finally {
                document.getElementById('integrated-loading').style.display = 'none';
            }
        });

        // Document parsing form submission
        document.getElementById('parsing-form').addEventListener('submit', async function(e) {
            e.preventDefault();

            const formData = new FormData(this);
            formData.set('overlap', document.getElementById('parsing-overlap').checked);
            formData.set('clean_for_ai', document.getElementById('parsing-clean-ai').checked);
            formData.set('replace_whitespace', document.getElementById('parsing-replace-whitespace').checked);
            formData.set('remove_urls_emails', document.getElementById('parsing-remove-urls-emails').checked);

            document.getElementById('parsing-loading').style.display = 'block';
            document.getElementById('parsing-results').style.display = 'none';

            try {
                const response = await fetch('/api/process-document', {
                    method: 'POST',
                    body: formData
                });

                const data = await response.json();

                if (!response.ok) {
                    throw new Error(data.error || 'An error occurred');
                }

                // Store result for download
                window.parsingResults = data;

                // Show results
                document.getElementById('parsing-segments-count').textContent = data.segments_count;

                // Render segments
                const segmentsContainer = document.getElementById('parsing-segments');
                segmentsContainer.innerHTML = '';

                data.segments.forEach((segment, index) => {
                    const segmentDiv = document.createElement('div');
                    segmentDiv.className = 'panel';

                    const header = document.createElement('div');
                    header.className = 'panel-header';
                    header.innerHTML = `Segment ${index + 1} <button type="button" class="toggle-btn">+</button>`;

                    const body = document.createElement('div');
                    body.className = 'panel-body';
                    body.innerHTML = `<pre>${segment.content}</pre>`;

                    segmentDiv.appendChild(header);
                    segmentDiv.appendChild(body);
                    segmentsContainer.appendChild(segmentDiv);

                    header.addEventListener('click', () => {
                        body.classList.toggle('active');
                        const btn = header.querySelector('.toggle-btn');
                        btn.textContent = body.classList.contains('active') ? '-' : '+';
                    });
                });

                document.getElementById('parsing-results').style.display = 'block';
                document.getElementById('dialogue-segments').value = JSON.stringify(data.segments);

            } catch (error) {
                console.error('Error:', error);
                const resultsDiv = document.getElementById('parsing-results');
                resultsDiv.style.display = 'block';
                resultsDiv.innerHTML = `<div class="error">Error: ${error.message}</div>`;
            } finally {
                document.getElementById('parsing-loading').style.display = 'none';
            }
        });

        // Dialogue generation form submission
        document.getElementById('dialogue-form').addEventListener('submit', async function(e) {
            e.preventDefault();

            let requestData;
            try {
                requestData = {
                    segments: JSON.parse(document.getElementById('dialogue-segments').value),
                    api_url: document.getElementById('dialogue-api-url').value,
                    api_key: document.getElementById('dialogue-api-key').value,
                    model: document.getElementById('dialogue-model').value,
                    dialogue_rounds: parseInt(document.getElementById('dialogue-rounds').value),
                    max_workers: parseInt(document.getElementById('dialogue-max-workers').value),
                    language: document.getElementById('dialogue-language').value,
                    max_generation_mode: document.getElementById('dialogue-max-generation').checked,
                    max_questions: parseInt(document.getElementById('dialogue-max-questions').value),
                    dialogue_token_limit: parseInt(document.getElementById('dialogue-token-limit').value)
                };
            } catch (error) {
                alert('Invalid segments JSON: ' + error.message);
                return;
            }

            document.getElementById('dialogue-loading').style.display = 'block';
            document.getElementById('dialogue-results').style.display = 'none';

            try {
                const response = await fetch('/api/generate-dialogues', {
                    method: 'POST',
                    headers: {
                        'Content-Type': 'application/json'
                    },
                    body: JSON.stringify(requestData)
                });

                const data = await response.json();

                if (!response.ok) {
                    throw new Error(data.error || 'An error occurred');
                }

                // Store result for download
                window.dialogueResults = data;

                // Show results
                document.getElementById('dialogue-dialogues-count').textContent = data.dialogues_count;

                // Render dialogues
                const dialoguesContainer = document.getElementById('dialogue-content');
                dialoguesContainer.innerHTML = '';

                data.dialogues.forEach((result, index) => {
                    const dialogueDiv = document.createElement('div');
                    dialogueDiv.className = 'panel';

                    const header = document.createElement('div');
                    header.className = 'panel-header';
                    header.innerHTML = `Segment ${index + 1} Dialogues <span class="badge">${Object.keys(result.dialogues).length}</span> <button type="button" class="toggle-btn">+</button>`;

                    const body = document.createElement('div');
                    body.className = 'panel-body';

                    // For each question in this chunk
                    Object.entries(result.dialogues).forEach(([qKey, qData]) => {
                        const questionDiv = document.createElement('div');
                        questionDiv.className = 'dialogue';

                        const questionHeader = document.createElement('h4');
                        questionHeader.textContent = qData.question;
                        questionDiv.appendChild(questionHeader);

                        // Display conversation
                        qData.dialogue.forEach(msg => {
                            if (msg.role === 'system') return; // Skip system messages

                            const messageDiv = document.createElement('div');
                            messageDiv.className = `message ${msg.role}`;
                            messageDiv.textContent = msg.content;
                            questionDiv.appendChild(messageDiv);
                        });

                        body.appendChild(questionDiv);
                    });

                    dialogueDiv.appendChild(header);
                    dialogueDiv.appendChild(body);
                    dialoguesContainer.appendChild(dialogueDiv);

                    header.addEventListener('click', () => {
                        body.classList.toggle('active');
                        const btn = header.querySelector('.toggle-btn');
                        btn.textContent = body.classList.contains('active') ? '-' : '+';
                    });
                });

                document.getElementById('dialogue-results').style.display = 'block';

            } catch (error) {
                console.error('Error:', error);
                const resultsDiv = document.getElementById('dialogue-results');
                resultsDiv.style.display = 'block';
                resultsDiv.innerHTML = `<div class="error">Error: ${error.message}</div>`;
            } finally {
                document.getElementById('dialogue-loading').style.display = 'none';
            }
        });

        // Download buttons
        document.getElementById('integrated-download-btn').addEventListener('click', async function() {
            if (!window.integratedResults) return;

            try {
                const response = await fetch('/api/download-results', {
                    method: 'POST',
                    headers: {
                        'Content-Type': 'application/json'
                    },
                    body: JSON.stringify(window.integratedResults)
                });

                if (!response.ok) {
                    throw new Error('Download failed');
                }

                const blob = await response.blob();
                const url = URL.createObjectURL(blob);
                const a = document.createElement('a');
                a.href = url;
                a.download = 'integrated_results.json';
                document.body.appendChild(a);
                a.click();
                URL.revokeObjectURL(url);
                document.body.removeChild(a);

            } catch (error) {
                console.error('Error:', error);
                alert('Download failed: ' + error.message);
            }
        });

        document.getElementById('parsing-download-btn').addEventListener('click', async function() {
            if (!window.parsingResults) return;

            try {
                const response = await fetch('/api/download-results', {
                    method: 'POST',
                    headers: {
                        'Content-Type': 'application/json'
                    },
                    body: JSON.stringify(window.parsingResults)
                });

                if (!response.ok) {
                    throw new Error('Download failed');
                }

                const blob = await response.blob();
                const url = URL.createObjectURL(blob);
                const a = document.createElement('a');
                a.href = url;
                a.download = 'document_segments.json';
                document.body.appendChild(a);
                a.click();
                URL.revokeObjectURL(url);
                document.body.removeChild(a);

            } catch (error) {
                console.error('Error:', error);
                alert('Download failed: ' + error.message);
            }
        });

        document.getElementById('dialogue-download-btn').addEventListener('click', async function() {
            if (!window.dialogueResults) return;

            try {
                const response = await fetch('/api/download-results', {
                    method: 'POST',
                    headers: {
                        'Content-Type': 'application/json'
                    },
                    body: JSON.stringify(window.dialogueResults)
                });

                if (!response.ok) {
                    throw new Error('Download failed');
                }

                const blob = await response.blob();
                const url = URL.createObjectURL(blob);
                const a = document.createElement('a');
                a.href = url;
                a.download = 'dialogue_results.json';
                document.body.appendChild(a);
                a.click();
                URL.revokeObjectURL(url);
                document.body.removeChild(a);

            } catch (error) {
                console.error('Error:', error);
                alert('Download failed: ' + error.message);
            }
        });
    </script>
</body>
</html>
            ''')

        # Start the web server
        port = int(os.environ.get('PORT', args.port))
        host = os.environ.get('HOST', args.host)
        debug = os.environ.get('DEBUG', str(args.debug)).lower() == 'true'

        logger.info(f"Starting server on {host}:{port} (debug={debug})")
        app.run(host=host, port=port, debug=debug)

    elif args.command == 'document':
        try:
            segments = process_document(
                args.file_path,
                languages=args.languages,
                segment_size=args.segment_size,
                overlap=not args.no_overlap,
                overlap_limit=args.overlap_limit,
                clean_for_ai=not args.no_clean_ai,
                segmentation_type=args.segmentation_type,
                replace_whitespace=args.replace_whitespace,
                remove_urls_emails=args.remove_urls_emails
            )

            print(f"Document processed into {len(segments)} segments")

            if args.output:
                with open(args.output, 'w', encoding='utf-8') as f:
                    json.dump(segments, f, ensure_ascii=False, indent=2)
                print(f"Results saved to {args.output}")
            else:
                # Print sample of first few segments
                for i, segment in enumerate(segments[:3]):
                    print(f"Segment {i + 1}:")
                    print(f"Type: {segment.get('type', 'unknown')}")
                    content_preview = segment['content'][:100] + "..." if len(segment['content']) > 100 else segment[
                        'content']
                    print(f"Content: {content_preview}")
                    print("-" * 50)

                if len(segments) > 3:
                    print(f"...and {len(segments) - 3} more segments")

        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            import traceback
            traceback.print_exc()

    elif args.command == 'dialogue':
        try:
            # Load input file with text chunks
            with open(args.input_file, 'r', encoding='utf-8') as f:
                data = json.load(f)

            # Extract text content
            if isinstance(data, list):
                chunks = [item.get('content', '') for item in data]
            else:
                if 'segments' in data:
                    chunks = [item.get('content', '') for item in data['segments']]
                else:
                    chunks = list(data.values())

            print(f"Loaded {len(chunks)} text chunks for dialogue generation")

            # Initialize dialogue generator
            generator = ParallelDialogueGenerator(
                api_url=args.api_url,
                api_key=args.api_key,
                output_dir=args.output_dir,
                model=args.model,
                max_generation_mode=args.max_generation_mode,
                max_questions=args.max_questions,
                dialogue_token_limit=args.token_limit,
                language=args.language
            )

            # Process chunks
            results = generator.process_chunks_parallel(
                chunks=chunks,
                rounds=args.rounds,
                max_workers=args.max_workers
            )

            print(f"Generated dialogues for {len(results)} chunks")

            # Save overall results
            overall_results_path = os.path.join(args.output_dir, 'all_dialogues.json')
            with open(overall_results_path, 'w', encoding='utf-8') as f:
                json.dump(results, f, ensure_ascii=False, indent=2)

            print(f"All dialogues saved to {overall_results_path}")

        except Exception as e:
            logger.error(f"Error generating dialogues: {str(e)}")
            import traceback
            traceback.print_exc()

    elif args.command == 'integrated':
        try:
            # Run integrated workflow
            result = process_document_and_generate_dialogues(
                file_path=args.file_path,
                api_url=args.api_url,
                api_key=args.api_key,
                model=args.model,
                dialogue_rounds=args.dialogue_rounds,
                max_workers=args.max_workers,
                output_dir=args.output_dir,
                language=args.language,
                max_generation_mode=args.max_generation_mode,
                max_questions=args.max_questions,
                dialogue_token_limit=args.token_limit,
                segment_size=args.segment_size,
                overlap=not args.no_overlap,
                overlap_limit=args.overlap_limit,
                clean_for_ai=not args.no_clean_ai,
                segmentation_type=args.segmentation_type,
                languages=args.languages
            )

            if result["status"] == "success":
                print(f"Integration workflow completed successfully:")
                print(f"- Document processed into {result['segments_count']} segments")
                print(f"- Generated dialogues for {result['dialogues_count']} segments")
                print(f"- Results saved to {result['results_file']}")
            else:
                print(f"Integration workflow failed: {result['error']}")

        except Exception as e:
            logger.error(f"Error in integrated workflow: {str(e)}")
            import traceback
            traceback.print_exc()

    else:
        parser.print_help()


if __name__ == "__main__":
    main()